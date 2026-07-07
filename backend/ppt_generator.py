import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from models import ESGRequest, ESGScores, AestheticTheme, PresentationType
from visual_kit import pillar_hero, icon_png, ring_png
from i18n import L


def _hexstr(rgb) -> str:
    """RGBColor -> '#RRGGBB' pour visual_kit."""
    return "#" + str(rgb)

# Color palettes per theme
THEMES = {
    AestheticTheme.CORPORATE_BLUE: {
        "bg_primary": RGBColor(0x1B, 0x3A, 0x6B),
        "bg_secondary": RGBColor(0xF4, 0xF7, 0xFD),
        "accent": RGBColor(0xF3, 0x9C, 0x12),
        "env": RGBColor(0x27, 0xAE, 0x60),
        "social": RGBColor(0x2E, 0x86, 0xC1),
        "gov": RGBColor(0x8E, 0x44, 0xAD),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark": RGBColor(0x2C, 0x3E, 0x50),
        "subtitle": RGBColor(0xAE, 0xC6, 0xE8),
        "muted": RGBColor(0x6B, 0x7C, 0x93),
        "card_bg": RGBColor(0xEB, 0xF2, 0xFB),
    },
    AestheticTheme.GREEN_NATURE: {
        "bg_primary": RGBColor(0x1A, 0x5C, 0x38),
        "bg_secondary": RGBColor(0xF3, 0xFB, 0xF5),
        "accent": RGBColor(0xF1, 0xC4, 0x0F),
        "env": RGBColor(0x2E, 0xCC, 0x71),
        "social": RGBColor(0x34, 0x98, 0xDB),
        "gov": RGBColor(0xE6, 0x7E, 0x22),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark": RGBColor(0x1A, 0x33, 0x25),
        "subtitle": RGBColor(0xA9, 0xDB, 0xBC),
        "muted": RGBColor(0x5E, 0x8C, 0x72),
        "card_bg": RGBColor(0xDD, 0xF3, 0xE4),
    },
    AestheticTheme.DARK_PREMIUM: {
        "bg_primary": RGBColor(0x0D, 0x11, 0x17),
        "bg_secondary": RGBColor(0x16, 0x1B, 0x22),
        "accent": RGBColor(0xF7, 0xC9, 0x48),
        "env": RGBColor(0x3F, 0xB9, 0x50),
        "social": RGBColor(0x58, 0xA6, 0xFF),
        "gov": RGBColor(0xBC, 0x8C, 0xFF),
        "text_light": RGBColor(0xE6, 0xED, 0xF3),
        "text_dark": RGBColor(0xE6, 0xED, 0xF3),
        "subtitle": RGBColor(0x8B, 0x94, 0x9E),
        "muted": RGBColor(0x8B, 0x94, 0x9E),
        "card_bg": RGBColor(0x1F, 0x25, 0x2E),
    },
    AestheticTheme.MINIMAL_WHITE: {
        "bg_primary": RGBColor(0x21, 0x21, 0x21),
        "bg_secondary": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0xFF, 0x6F, 0x00),
        "env": RGBColor(0x43, 0xA0, 0x47),
        "social": RGBColor(0x1E, 0x88, 0xE5),
        "gov": RGBColor(0x8E, 0x24, 0xAA),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark": RGBColor(0x21, 0x21, 0x21),
        "subtitle": RGBColor(0x9E, 0x9E, 0x9E),
        "muted": RGBColor(0x9E, 0x9E, 0x9E),
        "card_bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    AestheticTheme.SUNSET_TERRACOTTA: {
        "bg_primary": RGBColor(0x9A, 0x34, 0x12),
        "bg_secondary": RGBColor(0xFD, 0xF3, 0xEC),
        "accent": RGBColor(0xF4, 0xA2, 0x61),
        "env": RGBColor(0x2A, 0x9D, 0x8F),
        "social": RGBColor(0xE7, 0x6F, 0x51),
        "gov": RGBColor(0x6D, 0x59, 0x7A),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark": RGBColor(0x4A, 0x2C, 0x22),
        "subtitle": RGBColor(0xF2, 0xC9, 0xB0),
        "muted": RGBColor(0xA8, 0x70, 0x5F),
        "card_bg": RGBColor(0xFA, 0xE5, 0xD8),
    },
    AestheticTheme.OCEAN_DEEP: {
        "bg_primary": RGBColor(0x0F, 0x4C, 0x5C),
        "bg_secondary": RGBColor(0xEF, 0xF9, 0xFB),
        "accent": RGBColor(0x00, 0xBF, 0xA6),
        "env": RGBColor(0x43, 0xAA, 0x8B),
        "social": RGBColor(0x27, 0x7D, 0xA1),
        "gov": RGBColor(0x57, 0x75, 0x90),
        "text_light": RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark": RGBColor(0x12, 0x3B, 0x44),
        "subtitle": RGBColor(0xA8, 0xD8, 0xE0),
        "muted": RGBColor(0x5C, 0x8A, 0x96),
        "card_bg": RGBColor(0xDC, 0xF1, 0xF5),
    },
    AestheticTheme.ROYAL_PURPLE: {
        "bg_primary": RGBColor(0x2B, 0x10, 0x55),
        "bg_secondary": RGBColor(0x24, 0x10, 0x47),
        "accent": RGBColor(0xFF, 0xD5, 0x4F),
        "env": RGBColor(0x2E, 0x9E, 0x62),
        "social": RGBColor(0x7E, 0x9B, 0xF5),
        "gov": RGBColor(0xC0, 0x8C, 0xF5),
        "text_light": RGBColor(0xF2, 0xED, 0xFB),
        "text_dark": RGBColor(0xF2, 0xED, 0xFB),
        "subtitle": RGBColor(0xB9, 0xA6, 0xE0),
        "muted": RGBColor(0xB9, 0xA6, 0xE0),
        "card_bg": RGBColor(0x3A, 0x21, 0x70),
    },
}

# Design language per theme: fonts, header style, card style, cover layout
STYLES = {
    AestheticTheme.CORPORATE_BLUE: {
        "font_title": "Calibri", "font_body": "Calibri",
        "header": "band", "card": "flat", "cover": "classic",
        "dark_slides": False, "uppercase_titles": False,
    },
    AestheticTheme.GREEN_NATURE: {
        "font_title": "Trebuchet MS", "font_body": "Trebuchet MS",
        "header": "pill", "card": "rounded", "cover": "organic",
        "dark_slides": False, "uppercase_titles": False,
    },
    AestheticTheme.DARK_PREMIUM: {
        "font_title": "Georgia", "font_body": "Georgia",
        "header": "hairline", "card": "dark", "cover": "luxe",
        "dark_slides": True, "uppercase_titles": True,
    },
    AestheticTheme.MINIMAL_WHITE: {
        "font_title": "Segoe UI", "font_body": "Segoe UI",
        "header": "minimal", "card": "outline", "cover": "minimal",
        "dark_slides": False, "uppercase_titles": True,
    },
    AestheticTheme.SUNSET_TERRACOTTA: {
        "font_title": "Cambria", "font_body": "Calibri",
        "header": "pill", "card": "rounded", "cover": "organic",
        "dark_slides": False, "uppercase_titles": False,
    },
    AestheticTheme.OCEAN_DEEP: {
        "font_title": "Segoe UI", "font_body": "Segoe UI",
        "header": "band", "card": "flat", "cover": "classic",
        "dark_slides": False, "uppercase_titles": False,
    },
    AestheticTheme.ROYAL_PURPLE: {
        "font_title": "Georgia", "font_body": "Georgia",
        "header": "hairline", "card": "dark", "cover": "luxe",
        "dark_slides": True, "uppercase_titles": True,
    },
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

RECT = 1
ROUNDED_RECT = 5
OVAL = 9


def add_shape(slide, shape_type, left, top, width, height, fill: RGBColor = None,
              line_color: RGBColor = None, line_width_pt: float = None):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color is not None:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_bg_rect(slide, left, top, width, height, color: RGBColor):
    return add_shape(slide, RECT, left, top, width, height, fill=color)


def add_text(slide, text, left, top, width, height, font_size=18,
             bold=False, color: RGBColor = RGBColor(0, 0, 0),
             align=PP_ALIGN.LEFT, italic=False, word_wrap=True,
             font: str = None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if font:
        run.font.name = font
    return txBox


def add_image_from_bytes(slide, img_bytes, left, top, width=None, height=None):
    img_io = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_io, left, top, width, height)


def maybe_upper(text: str, style: dict) -> str:
    return text.upper() if style["uppercase_titles"] else text


def _fit_title(base, title):
    """Réduit la taille du titre s'il est long (titres-conclusion)."""
    n = len(title)
    if n > 52:
        return base * 0.66
    if n > 38:
        return base * 0.78
    if n > 28:
        return base * 0.9
    return base


def content_slide(prs, blank_layout, theme, style, title, color: RGBColor, kicker=None):
    """Slide à en-tête thématisé. `kicker` = petit sur-titre optionnel."""
    slide = prs.slides.add_slide(blank_layout)
    title = maybe_upper(title, style)
    ft, fb = style["font_title"], style["font_body"]

    if style["header"] == "band":
        add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_secondary"])
        add_bg_rect(slide, 0, 0, SLIDE_W, Inches(1.25), color)
        add_bg_rect(slide, 0, Inches(1.25), Inches(0.07), SLIDE_H - Inches(1.25), theme["accent"])
        if kicker:
            add_text(slide, kicker.upper(), Inches(0.32), Inches(0.14), Inches(12), Inches(0.32),
                     font_size=11, bold=True, color=theme["accent"], font=fb)
        add_text(slide, title, Inches(0.3), Inches(0.42 if kicker else 0.2), Inches(12.7), Inches(0.78),
                 font_size=_fit_title(26, title), bold=True, color=theme["text_light"], font=ft)

    elif style["header"] == "pill":
        add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_secondary"])
        add_shape(slide, OVAL, Inches(11.8), Inches(-1.2), Inches(3), Inches(3), fill=theme["card_bg"])
        add_shape(slide, OVAL, Inches(12.6), Inches(6.4), Inches(2), Inches(2), fill=theme["card_bg"])
        add_shape(slide, ROUNDED_RECT, Inches(0.3), Inches(0.22), Inches(11.2), Inches(0.95), fill=color)
        add_text(slide, title, Inches(0.6), Inches(0.3), Inches(10.8), Inches(0.8),
                 font_size=_fit_title(22, title), bold=True, color=theme["text_light"], font=ft)

    elif style["header"] == "hairline":
        add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_primary"])
        if kicker:
            add_text(slide, kicker.upper(), Inches(0.5), Inches(0.24), Inches(12), Inches(0.3),
                     font_size=11, bold=True, color=theme["accent"], font=fb)
        add_text(slide, title, Inches(0.5), Inches(0.5 if kicker else 0.32), Inches(12.3), Inches(0.7),
                 font_size=_fit_title(24, title), bold=False, color=color, font=ft)
        add_bg_rect(slide, Inches(0.5), Inches(1.28), Inches(12.3), Inches(0.02), theme["accent"])

    else:  # minimal
        add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_secondary"])
        add_bg_rect(slide, Inches(0.5), Inches(0.42), Inches(0.22), Inches(0.22), color)
        add_text(slide, title, Inches(0.9), Inches(0.26), Inches(11.8), Inches(0.75),
                 font_size=_fit_title(20, title), bold=True, color=theme["text_dark"], font=ft)
        add_bg_rect(slide, Inches(0.5), Inches(1.12), Inches(2.2), Inches(0.02), theme["text_dark"])

    return slide


def kpi_card(slide, left, top, w, h, theme, style, color: RGBColor):
    """Draw a KPI card background in the theme's card style."""
    if style["card"] == "flat":
        add_bg_rect(slide, left, top, w, h, theme["card_bg"])
        add_bg_rect(slide, left, top, w, Inches(0.1), color)
    elif style["card"] == "rounded":
        add_shape(slide, ROUNDED_RECT, left, top, w, h, fill=theme["card_bg"])
        add_shape(slide, OVAL, left + Inches(0.15), top + Inches(0.15),
                  Inches(0.18), Inches(0.18), fill=color)
    elif style["card"] == "dark":
        add_bg_rect(slide, left, top, w, h, theme["card_bg"])
        add_bg_rect(slide, left, top, w, Inches(0.03), theme["accent"])
    else:  # outline
        add_shape(slide, RECT, left, top, w, h, fill=theme["card_bg"],
                  line_color=RGBColor(0xE0, 0xE0, 0xE0), line_width_pt=1.0)
        add_bg_rect(slide, left, top, Inches(0.06), h, color)


def kpi_grid(slide, kpis, theme, style, color: RGBColor, top_start=Inches(1.3), ncols=3):
    fb = style["font_body"]
    cols = [Inches(0.3), Inches(4.55), Inches(8.8)][:ncols]
    max_cards = ncols * 3
    label_indent = Inches(0.42) if style["card"] == "rounded" else Inches(0.2)
    for idx, (kpi_label, kpi_val) in enumerate(kpis[:max_cards]):
        col = cols[idx % ncols]
        row = top_start + (idx // ncols) * Inches(1.9)
        kpi_card(slide, col, row, Inches(4.0), Inches(1.7), theme, style, color)
        add_text(slide, kpi_label, col + label_indent, row + Inches(0.15),
                 Inches(3.5), Inches(0.45), font_size=11, color=theme["muted"], font=fb)
        add_text(slide, kpi_val, col + Inches(0.2), row + Inches(0.65),
                 Inches(3.7), Inches(0.85), font_size=22, bold=True, color=color, font=fb)


# ── Covers ────────────────────────────────────────────────────────────────

def _cover_tag(TR, scores):
    b = "high" if scores.total_esg_score >= 75 else "good" if scores.total_esg_score >= 60 \
        else "mid" if scores.total_esg_score >= 45 else "low"
    return TR["cover_tag_" + b]


def cover_classic(slide, theme, style, request, scores, subtitle, TR):
    ft, fb = style["font_title"], style["font_body"]
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_primary"])
    add_bg_rect(slide, 0, 0, Inches(0.18), SLIDE_H, theme["accent"])
    # Aplat sombre en bas pour ancrer la composition (magazine)
    add_bg_rect(slide, 0, Inches(5.4), SLIDE_W, Inches(2.1), theme["bg_primary"])

    # Kicker + nom géant + accroche (bloc éditorial gauche)
    add_text(slide, subtitle.upper(), Inches(0.65), Inches(1.15), Inches(9), Inches(0.5),
             font_size=15, bold=True, color=theme["accent"], font=fb)
    add_text(slide, request.company.name.upper(), Inches(0.6), Inches(1.75), Inches(9.2), Inches(1.9),
             font_size=52, bold=True, color=theme["text_light"], font=ft)
    add_bg_rect(slide, Inches(0.65), Inches(3.75), Inches(2.0), Inches(0.06), theme["accent"])
    add_text(slide, _cover_tag(TR, scores), Inches(0.65), Inches(4.0), Inches(8.5), Inches(0.7),
             font_size=21, italic=True, color=theme["subtitle"], font=fb)

    # Bandeau bas : méta à gauche, score mis en scène à droite
    add_text(slide, f"{TR['exercise']} {request.company.reporting_year}   |   {request.company.sector}   |   {request.company.country}",
             Inches(0.65), Inches(6.55), Inches(8), Inches(0.5),
             font_size=13, color=theme["subtitle"], font=fb)
    add_text(slide, f"{scores.total_esg_score:.0f}", Inches(8.55), Inches(5.35), Inches(2.4), Inches(1.6),
             font_size=82, bold=True, color=theme["accent"], align=PP_ALIGN.RIGHT, font=ft)
    add_text(slide, "/100", Inches(8.55), Inches(6.75), Inches(2.4), Inches(0.4),
             font_size=14, color=theme["subtitle"], align=PP_ALIGN.RIGHT, font=fb)
    add_shape(slide, ROUNDED_RECT, Inches(11.2), Inches(5.55), Inches(1.55), Inches(0.95), fill=theme["accent"])
    add_text(slide, scores.rating, Inches(11.2), Inches(5.68), Inches(1.55), Inches(0.7),
             font_size=32, bold=True, color=theme["bg_primary"], align=PP_ALIGN.CENTER, font=ft)


def cover_organic(slide, theme, style, request, scores, subtitle, TR):
    ft, fb = style["font_title"], style["font_body"]
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_secondary"])
    # Cercles organiques (identité) en haut à droite
    add_shape(slide, OVAL, Inches(9.6), Inches(-2.4), Inches(6.0), Inches(6.0), fill=theme["card_bg"])
    add_shape(slide, OVAL, Inches(11.3), Inches(-0.9), Inches(3.6), Inches(3.6), fill=theme["env"])
    add_shape(slide, OVAL, Inches(10.5), Inches(1.7), Inches(1.3), Inches(1.3), fill=theme["accent"])
    # Bloc éditorial
    add_text(slide, subtitle.upper(), Inches(0.75), Inches(1.35), Inches(8.5), Inches(0.5),
             font_size=15, bold=True, color=theme["env"], font=fb)
    add_text(slide, request.company.name, Inches(0.7), Inches(1.95), Inches(9), Inches(1.7),
             font_size=48, bold=True, color=theme["bg_primary"], font=ft)
    add_shape(slide, ROUNDED_RECT, Inches(0.75), Inches(3.75), Inches(2.0), Inches(0.08), fill=theme["accent"])
    add_text(slide, _cover_tag(TR, scores), Inches(0.75), Inches(4.0), Inches(8.5), Inches(0.7),
             font_size=21, italic=True, color=theme["muted"], font=fb)
    # Bandeau bas : score mis en scène
    add_shape(slide, ROUNDED_RECT, Inches(0.7), Inches(5.35), Inches(11.93), Inches(1.5), fill=theme["bg_primary"])
    add_text(slide, f"{TR['exercise']} {request.company.reporting_year}   |   {request.company.sector}   |   {request.company.country}",
             Inches(1.15), Inches(5.65), Inches(7), Inches(0.5), font_size=13, color=theme["subtitle"], font=fb)
    add_text(slide, f"{TR['chart_global']}", Inches(1.15), Inches(6.15), Inches(7), Inches(0.5),
             font_size=13, bold=True, color=theme["text_light"], font=fb)
    add_text(slide, f"{scores.total_esg_score:.0f}", Inches(8.7), Inches(5.42), Inches(2.2), Inches(1.4),
             font_size=54, bold=True, color=theme["accent"], align=PP_ALIGN.RIGHT, font=ft)
    add_shape(slide, OVAL, Inches(11.15), Inches(5.62), Inches(1.0), Inches(1.0), fill=theme["accent"])
    add_text(slide, scores.rating, Inches(11.15), Inches(5.78), Inches(1.0), Inches(0.6),
             font_size=24, bold=True, color=theme["bg_primary"], align=PP_ALIGN.CENTER, font=ft)


def cover_luxe(slide, theme, style, request, scores, subtitle, TR):
    ft, fb = style["font_title"], style["font_body"]
    gold = theme["accent"]
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_primary"])
    # Fin cadre or (élégance)
    m, tk = Inches(0.4), Inches(0.015)
    add_bg_rect(slide, m, m, SLIDE_W - m - m, tk, gold)
    add_bg_rect(slide, m, SLIDE_H - m, SLIDE_W - m - m, tk, gold)
    add_bg_rect(slide, m, m, tk, SLIDE_H - m - m, gold)
    add_bg_rect(slide, SLIDE_W - m, m, tk, SLIDE_H - m - m + tk, gold)
    # Composition serif centrée
    add_text(slide, subtitle.upper(), Inches(1), Inches(1.15), Inches(11.33), Inches(0.5),
             font_size=14, bold=True, color=gold, align=PP_ALIGN.CENTER, font=fb)
    add_text(slide, request.company.name.upper(), Inches(1), Inches(1.95), Inches(11.33), Inches(1.4),
             font_size=46, bold=False, color=theme["text_light"], align=PP_ALIGN.CENTER, font=ft)
    add_bg_rect(slide, Inches(6.06), Inches(3.5), Inches(1.2), Inches(0.02), gold)
    add_text(slide, _cover_tag(TR, scores), Inches(1), Inches(3.75), Inches(11.33), Inches(0.6),
             font_size=17, italic=True, color=theme["subtitle"], align=PP_ALIGN.CENTER, font=ft)
    # Note mise en scène (le grand signe de qualité)
    add_text(slide, scores.rating, Inches(1), Inches(4.6), Inches(11.33), Inches(1.3),
             font_size=64, bold=True, color=gold, align=PP_ALIGN.CENTER, font=ft)
    add_text(slide, f"{TR['chart_global']}  {scores.total_esg_score:.0f} / 100", Inches(1), Inches(6.0),
             Inches(11.33), Inches(0.5), font_size=15, color=theme["subtitle"], align=PP_ALIGN.CENTER, font=fb)
    add_text(slide, f"{TR['exercise']} {request.company.reporting_year}  •  {request.company.sector}  •  {request.company.country}",
             Inches(1), Inches(6.55), Inches(11.33), Inches(0.4),
             font_size=12, italic=True, color=theme["subtitle"], align=PP_ALIGN.CENTER, font=ft)


def cover_minimal(slide, theme, style, request, scores, subtitle, TR):
    ft, fb = style["font_title"], style["font_body"]
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_secondary"])
    add_bg_rect(slide, Inches(0.7), Inches(1.05), Inches(0.4), Inches(0.4), theme["accent"])
    add_text(slide, subtitle.upper(), Inches(0.72), Inches(1.65), Inches(11), Inches(0.5),
             font_size=13, bold=True, color=theme["muted"], font=fb)
    add_text(slide, request.company.name, Inches(0.65), Inches(2.25), Inches(11.7), Inches(1.7),
             font_size=54, bold=True, color=theme["text_dark"], font=ft)
    add_text(slide, _cover_tag(TR, scores), Inches(0.7), Inches(4.0), Inches(9), Inches(0.6),
             font_size=19, italic=True, color=theme["muted"], font=fb)
    add_bg_rect(slide, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.02), RGBColor(0xBD, 0xBD, 0xBD))
    # Méta gauche + score mis en scène droite
    add_text(slide, f"{TR['exercise']} {request.company.reporting_year}", Inches(0.7), Inches(5.15),
             Inches(4), Inches(0.4), font_size=13, color=theme["muted"], font=fb)
    add_text(slide, request.company.sector, Inches(0.7), Inches(5.55), Inches(5), Inches(0.4),
             font_size=13, color=theme["muted"], font=fb)
    add_text(slide, request.company.country, Inches(0.7), Inches(5.95), Inches(5), Inches(0.4),
             font_size=13, color=theme["muted"], font=fb)
    add_text(slide, f"{scores.total_esg_score:.0f}", Inches(8.3), Inches(5.0), Inches(2.6), Inches(1.6),
             font_size=88, bold=True, color=theme["accent"], align=PP_ALIGN.RIGHT, font=ft)
    add_text(slide, f"/100  —  {TR['note']} {scores.rating}", Inches(6.9), Inches(6.5), Inches(4.0), Inches(0.4),
             font_size=15, bold=True, color=theme["text_dark"], align=PP_ALIGN.RIGHT, font=fb)


COVERS = {"classic": cover_classic, "organic": cover_organic,
          "luxe": cover_luxe, "minimal": cover_minimal}


def pillar_infographic(prs, blank_layout, theme, style, pillar_key,
                       title, subtitle, score, kpis, t, insight=""):
    """Slide de pilier en infographie : héros illustré à gauche + chips KPI.

    kpis : liste de (icon_name, value_str, label). 6 max affichés.
    """
    ft, fb = style["font_title"], style["font_body"]
    color = theme[pillar_key]
    dark = style["dark_slides"]
    white = RGBColor(0xFF, 0xFF, 0xFF)

    slide = prs.slides.add_slide(blank_layout)
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_primary"] if dark else theme["bg_secondary"])

    # ── Héros illustré (flanc gauche) ──────────────────────────────
    HERO_W = Inches(4.5)
    pal = {"top": _hexstr(theme["bg_primary"]), "bottom": _hexstr(color),
           "accent": _hexstr(theme["accent"]), "light": _hexstr(theme["card_bg"])}
    try:
        hero = pillar_hero(pillar_key, pal, 470, 780)
        add_image_from_bytes(slide, hero, 0, 0, HERO_W, SLIDE_H)
    except Exception:
        add_bg_rect(slide, 0, 0, HERO_W, SLIDE_H, color)

    # Kicker = identité du pilier ; grand titre = la conclusion (information scent)
    add_text(slide, title.upper(), Inches(0.37), Inches(0.42), Inches(3.85), Inches(0.35),
             font_size=13, bold=True, color=theme["accent"], font=fb)
    add_text(slide, subtitle, Inches(0.35), Inches(0.85), Inches(3.95), Inches(2.5),
             font_size=_fit_title(24, subtitle), bold=True, color=white, font=ft)

    # Anneau de score en bas du héros
    try:
        ring = ring_png(score, 320, _hexstr(theme["accent"]))
        add_image_from_bytes(slide, ring, Inches(1.4), Inches(4.75), Inches(1.75), Inches(1.75))
    except Exception:
        pass
    add_text(slide, f"{score:.0f}", Inches(1.4), Inches(5.28), Inches(1.75), Inches(0.7),
             font_size=38, bold=True, color=white, align=PP_ALIGN.CENTER, font=ft)
    add_text(slide, t["score_100_caps"], Inches(1.0), Inches(6.65), Inches(2.55), Inches(0.4),
             font_size=11, bold=True, color=white, align=PP_ALIGN.CENTER, font=fb)

    # ── En-tête droite ─────────────────────────────────────────────
    RX = Inches(4.95)
    RW = Inches(8.05)
    add_text(slide, t["key_indicators"].upper(), RX, Inches(0.5), Inches(8), Inches(0.5),
             font_size=15, bold=True, color=theme["muted"], font=fb)
    add_bg_rect(slide, RX, Inches(1.02), Inches(1.4), Inches(0.045), color)

    # ── Chips KPI (2 × 2, plus grandes et lisibles) ────────────────
    chip_w, chip_h = Inches(3.9), Inches(1.42)
    gap_x, gap_y = Inches(0.25), Inches(0.24)
    top0 = Inches(1.35)
    border = RGBColor(0x3A, 0x40, 0x4A) if dark else RGBColor(0xE2, 0xE8, 0xF0)
    circ = Inches(0.9)
    for i, (icon_name, value, label) in enumerate(kpis[:4]):
        col, row = i % 2, i // 2
        x = RX + col * (chip_w + gap_x)
        y = top0 + row * (chip_h + gap_y)
        add_shape(slide, ROUNDED_RECT, x, y, chip_w, chip_h, fill=theme["card_bg"],
                  line_color=border, line_width_pt=0.75)
        add_shape(slide, RECT, x, y + Inches(0.18), Inches(0.06), chip_h - Inches(0.36), fill=color)
        add_shape(slide, OVAL, x + Inches(0.24), y + Inches(0.27), circ, circ, fill=color)
        try:
            add_image_from_bytes(slide, icon_png(icon_name, 130, "#FFFFFF"),
                                 x + Inches(0.43), y + Inches(0.46), Inches(0.52), Inches(0.52))
        except Exception:
            pass
        add_text(slide, value, x + Inches(1.35), y + Inches(0.24), chip_w - Inches(1.45), Inches(0.65),
                 font_size=23, bold=True, color=color, font=ft)
        add_text(slide, label, x + Inches(1.36), y + Inches(0.9), chip_w - Inches(1.45), Inches(0.45),
                 font_size=10.5, color=theme["muted"], font=fb)

    # ── Carte « Lecture métier » (insight rédigé) ──────────────────
    iy = top0 + 2 * (chip_h + gap_y) + Inches(0.18)
    ih = Inches(2.02)
    add_shape(slide, ROUNDED_RECT, RX, iy, RW, ih,
              fill=color if dark else theme["card_bg"],
              line_color=color, line_width_pt=1.5)
    add_shape(slide, RECT, RX, iy, Inches(0.11), ih, fill=color)
    lbl_color = white if dark else color
    add_text(slide, t["key_takeaway"], RX + Inches(0.42), iy + Inches(0.24), RW - Inches(0.7), Inches(0.4),
             font_size=12.5, bold=True, color=lbl_color, font=fb)
    txt_color = white if dark else theme["text_dark"]
    add_text(slide, insight or "", RX + Inches(0.42), iy + Inches(0.72), RW - Inches(0.8), ih - Inches(0.85),
             font_size=16, color=txt_color, font=fb)
    return slide


def generate_pptx(request: ESGRequest, scores: ESGScores, content: dict,
                  chart_images: dict, logo_bytes: bytes = None) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    theme = THEMES.get(request.aesthetic_theme, THEMES[AestheticTheme.CORPORATE_BLUE])
    style = STYLES.get(request.aesthetic_theme, STYLES[AestheticTheme.CORPORATE_BLUE])
    ft, fb = style["font_title"], style["font_body"]
    ptype = request.presentation_type
    t = L(request.language)
    from content_generator import pillar_insights, score_verdict
    _insights = pillar_insights(request, scores)
    _verdict = score_verdict(request, scores)
    from content_generator import pillar_headline, section_headlines, benchmark_verdict, maturity_text
    _headlines = pillar_headline(request, scores)
    _sh = section_headlines(request, scores)
    _bv = benchmark_verdict(request, scores)
    _mat = maturity_text(request, scores)
    blank_layout = prs.slide_layouts[6]

    # ── SLIDE 1: Cover ──────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank_layout)
    subtitle_map = {
        PresentationType.EXECUTIVE_SUMMARY: t["pres_executive_summary"],
        PresentationType.INVESTOR_DECK: t["pres_investor_deck"],
        PresentationType.DETAILED_REPORT: t["pres_detailed_report"],
        PresentationType.STAKEHOLDER_BRIEF: t["pres_stakeholder_brief"],
        PresentationType.ANNUAL_REPORT: t["pres_annual_report"],
    }
    COVERS[style["cover"]](slide, theme, style, request, scores,
                           subtitle_map.get(ptype, t["pres_default"]), t)

    # Logo entreprise sur la couverture
    if logo_bytes:
        logo_pos = {
            "classic": (Inches(11.75), Inches(0.35)),
            "organic": (Inches(0.7), Inches(0.45)),
            "luxe": (Inches(6.17), Inches(0.55)),
            "minimal": (Inches(11.55), Inches(0.9)),
        }[style["cover"]]
        try:
            slide.shapes.add_picture(io.BytesIO(logo_bytes),
                                     logo_pos[0], logo_pos[1], height=Inches(1.0))
        except Exception:
            pass

    # Présentateur
    if request.company.presenter_name:
        line = f"{t['presented_by']} {request.company.presenter_name}"
        if request.company.presenter_title:
            line += f" — {request.company.presenter_title}"
        pres_y = Inches(6.9) if style["cover"] == "organic" else Inches(6.55)
        pres_align = PP_ALIGN.CENTER if style["cover"] == "luxe" else PP_ALIGN.LEFT
        pres_color = theme["muted"] if style["cover"] == "minimal" else theme["subtitle"]
        add_text(slide, line, Inches(0.7), pres_y, Inches(11.9), Inches(0.4),
                 font_size=12, italic=True, color=pres_color, font=fb, align=pres_align)

    # ── SLIDE 1bis: Illustration de couverture (générée localement) ──────
    if "cover_art" in chart_images:
        slide = prs.slides.add_slide(blank_layout)
        bg = theme["bg_primary"] if style["dark_slides"] else theme["bg_secondary"]
        add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, bg)
        add_image_from_bytes(slide, chart_images["cover_art"], 0, 0, SLIDE_W, Inches(4.4))
        add_text(slide, maybe_upper(f"{t['pres_default']} {request.company.reporting_year}", style),
                 Inches(0.7), Inches(4.9), Inches(11.9), Inches(0.9),
                 font_size=32, bold=True, color=theme["text_dark"], font=ft)
        add_text(slide, f"{request.company.name}  •  {request.company.sector}",
                 Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.5),
                 font_size=15, color=theme["muted"], font=fb)

    # ── SLIDE 2: Tableau de Bord ESG ────────────────────────────────────
    dark = style["dark_slides"]
    white = RGBColor(0xFF, 0xFF, 0xFF)
    slide = prs.slides.add_slide(blank_layout)
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_primary"] if dark else theme["bg_secondary"])

    # ── HERO gauche : score global mis en scène (élément dominant) ──
    HP = Inches(4.75)
    add_bg_rect(slide, 0, 0, HP, SLIDE_H, theme["bg_primary"])
    add_bg_rect(slide, HP - Inches(0.08), 0, Inches(0.08), SLIDE_H, theme["accent"])
    add_text(slide, t["exec_kicker"].format(y=request.company.reporting_year),
             Inches(0.55), Inches(0.75), Inches(3.9), Inches(0.5),
             font_size=13, bold=True, color=theme["accent"], font=fb)
    add_text(slide, f"{scores.total_esg_score:.0f}", Inches(0.2), Inches(1.85), Inches(4.35), Inches(2.1),
             font_size=125, bold=True, color=white, align=PP_ALIGN.CENTER, font=ft)
    add_text(slide, "/ 100", Inches(0.2), Inches(3.95), Inches(4.35), Inches(0.55),
             font_size=22, color=theme["subtitle"], align=PP_ALIGN.CENTER, font=fb)
    # badge note
    add_shape(slide, ROUNDED_RECT, Inches(1.55), Inches(4.75), Inches(1.65), Inches(0.85), fill=theme["accent"])
    add_text(slide, scores.rating, Inches(1.55), Inches(4.9), Inches(1.65), Inches(0.6),
             font_size=30, bold=True, color=theme["bg_primary"], align=PP_ALIGN.CENTER, font=ft)
    add_text(slide, t["global_score_caption"], Inches(0.2), Inches(5.85), Inches(4.35), Inches(0.4),
             font_size=12, bold=True, color=theme["subtitle"], align=PP_ALIGN.CENTER, font=fb)

    # ── Droite : titre éditorial + verdict (secondaire) ─────────────
    RX = Inches(5.2)
    add_text(slide, t["exec_title"], RX, Inches(0.65), Inches(7.9), Inches(0.9),
             font_size=34, bold=True, color=theme["text_dark"], font=ft)
    add_bg_rect(slide, RX + Inches(0.03), Inches(1.62), Inches(1.6), Inches(0.06), theme["accent"])
    add_text(slide, _verdict, RX, Inches(1.9), Inches(7.85), Inches(1.5),
             font_size=16, color=theme["text_dark"] if not dark else theme["subtitle"], font=fb)

    # ── Détail : barres piliers (gauche) + radar (droite) ───────────
    add_text(slide, t["pillars_caption"], RX, Inches(3.5), Inches(4), Inches(0.4),
             font_size=12, bold=True, color=theme["muted"], font=fb)
    bars = [(t["chart_env"], scores.environmental_score, theme["env"]),
            (t["chart_soc"], scores.social_score, theme["social"]),
            (t["chart_gov"], scores.governance_score, theme["gov"])]
    track = RGBColor(0x3A, 0x40, 0x4A) if dark else RGBColor(0xE2, 0xE8, 0xF0)
    bar_w = Inches(3.4)
    for i, (label, sc, col) in enumerate(bars):
        y = Inches(4.05) + i * Inches(0.92)
        add_text(slide, label, RX, y, Inches(3.4), Inches(0.32),
                 font_size=12.5, bold=True, color=theme["text_dark"], font=fb)
        add_text(slide, f"{sc:.0f}", RX + bar_w - Inches(0.6), y, Inches(0.6), Inches(0.32),
                 font_size=13, bold=True, color=col, align=PP_ALIGN.RIGHT, font=ft)
        add_shape(slide, ROUNDED_RECT, RX, y + Inches(0.36), bar_w, Inches(0.17), fill=track)
        fill_w = max(Inches(0.17), Inches(3.4) * max(0, min(100, sc)) / 100)
        add_shape(slide, ROUNDED_RECT, RX, y + Inches(0.36), fill_w, Inches(0.17), fill=col)

    if "radar" in chart_images:
        add_image_from_bytes(slide, chart_images["radar"],
                             Inches(9.15), Inches(3.55), height=Inches(3.55))

    header_color = theme["accent"] if style["header"] == "hairline" else theme["bg_primary"]

    # ── SLIDE 2b: Benchmark sectoriel + maturité (diagnostic) ───────────
    if "benchmark" in chart_images:
        _dk = style["dark_slides"]
        slide = content_slide(prs, blank_layout, theme, style,
                              _bv["title"], header_color, kicker=t["benchmark_kicker"])
        add_image_from_bytes(slide, chart_images["benchmark"],
                             Inches(0.35), Inches(1.7), width=Inches(7.5))
        # Carte lecture métier (positionnement)
        cx, cw = Inches(8.15), Inches(4.85)
        add_shape(slide, ROUNDED_RECT, cx, Inches(1.75), cw, Inches(2.35),
                  fill=theme["card_bg"], line_color=theme["accent"], line_width_pt=1.25)
        add_shape(slide, RECT, cx, Inches(1.75), Inches(0.09), Inches(2.35), fill=theme["accent"])
        add_text(slide, t["key_takeaway"], cx + Inches(0.3), Inches(1.93), cw - Inches(0.5), Inches(0.35),
                 font_size=12, bold=True, color=theme["accent"] if not _dk else theme["text_dark"], font=fb)
        add_text(slide, _bv["insight"], cx + Inches(0.3), Inches(2.35), cw - Inches(0.55), Inches(1.7),
                 font_size=14, color=theme["text_dark"], font=fb)
        # Échelle de maturité ESG (5 stades)
        add_text(slide, t["maturity_title"], cx + Inches(0.02), Inches(4.4), cw, Inches(0.35),
                 font_size=12, bold=True, color=theme["muted"], font=fb)
        stages = [t["mat_initiated"], t["mat_structuring"], t["mat_structured"], t["mat_advanced"], t["mat_exemplary"]]
        seg_w = Inches(0.92); seg_h = Inches(0.5); gap = Inches(0.06)
        for i, st in enumerate(stages):
            sx = cx + i * (seg_w + gap)
            active = (i == _mat["stage"])
            fill = theme["accent"] if active else theme["card_bg"]
            add_shape(slide, ROUNDED_RECT, sx, Inches(4.85), seg_w, seg_h,
                      fill=fill, line_color=theme["accent"] if not active else fill, line_width_pt=1.0)
            add_text(slide, str(i + 1), sx, Inches(4.9), seg_w, Inches(0.4),
                     font_size=14, bold=True, align=PP_ALIGN.CENTER,
                     color=(theme["bg_primary"] if active else theme["muted"]), font=ft)
        add_text(slide, stages[_mat["stage"]], cx, Inches(5.5), cw, Inches(0.4),
                 font_size=16, bold=True, color=theme["accent"] if not _dk else theme["text_dark"], font=ft)
        add_text(slide, _mat["next_hint"], cx, Inches(5.95), cw, Inches(0.7),
                 font_size=12.5, color=theme["text_dark"], font=fb)
        add_text(slide, t["bench_note"], Inches(0.4), Inches(6.95), Inches(7.5), Inches(0.35),
                 font_size=9, italic=True, color=theme["muted"], font=fb)

    # ── SLIDE 3: Environnement (infographie) ─────────────────────────────
    env = request.environmental
    env_kpis = []
    if env.co2_emissions_tonnes is not None:
        env_kpis.append(("cloud", f"{env.co2_emissions_tonnes:,.0f} t", t["kpi"]["co2_total"]))
    if env.renewable_energy_percent is not None:
        env_kpis.append(("recycle", f"{env.renewable_energy_percent:.0f}%", t["kpi"]["renewable"]))
    if env.energy_consumption_mwh is not None:
        env_kpis.append(("bolt", f"{env.energy_consumption_mwh:,.0f}", t["kpi"]["energy"]))
    if env.water_consumption_m3 is not None:
        env_kpis.append(("drop", f"{env.water_consumption_m3:,.0f}", t["kpi"]["water"]))
    if env.waste_recycled_percent is not None:
        env_kpis.append(("leaf", f"{env.waste_recycled_percent:.0f}%", t["kpi"]["recycling"]))
    if env.biodiversity_initiatives is not None:
        env_kpis.append(("tree", f"{env.biodiversity_initiatives}", t["kpi"]["biodiversity"]))
    if env.scope3_emissions is not None:
        env_kpis.append(("cloud", f"{env.scope3_emissions:,.0f} t", t["kpi"]["scope3"]))
    if env.waste_generated_tonnes is not None:
        env_kpis.append(("recycle", f"{env.waste_generated_tonnes:,.0f} t", t["kpi"]["waste"]))
    pillar_infographic(prs, blank_layout, theme, style, "env",
                       t["pillar_env"], _headlines["env"],
                       scores.environmental_score, env_kpis, t, _insights["env"])

    # ── SLIDE 4: Social (infographie) ────────────────────────────────────
    soc = request.social
    soc_kpis = []
    if soc.total_employees is not None:
        soc_kpis.append(("people", f"{soc.total_employees:,}", t["kpi"]["employees"]))
    if soc.female_employees_percent is not None:
        soc_kpis.append(("people", f"{soc.female_employees_percent:.0f}%", t["kpi"]["women_workforce"]))
    if soc.training_hours_per_employee is not None:
        soc_kpis.append(("cap", f"{soc.training_hours_per_employee:.0f} h", t["kpi"]["training"]))
    if soc.accident_frequency_rate is not None:
        soc_kpis.append(("shield", f"{soc.accident_frequency_rate:.1f}", t["kpi"]["accident_rate"]))
    if soc.customer_satisfaction_score is not None:
        soc_kpis.append(("heart", f"{soc.customer_satisfaction_score:.1f}/10", t["kpi"]["satisfaction"]))
    if soc.employee_turnover_percent is not None:
        soc_kpis.append(("chart", f"{soc.employee_turnover_percent:.0f}%", t["kpi"]["turnover"]))
    if soc.disabled_employees_percent is not None:
        soc_kpis.append(("heart", f"{soc.disabled_employees_percent:.1f}%", t["kpi"]["disabled"]))
    if soc.community_investment_eur is not None:
        soc_kpis.append(("heart", f"{soc.community_investment_eur:,.0f} €", t["kpi"]["community"]))
    pillar_infographic(prs, blank_layout, theme, style, "social",
                       t["pillar_soc"], _headlines["social"],
                       scores.social_score, soc_kpis, t, _insights["social"])

    # ── SLIDE 5: Gouvernance (infographie) ───────────────────────────────
    gov = request.governance
    gov_kpis = []
    if gov.board_members is not None:
        gov_kpis.append(("columns", f"{gov.board_members}", t["kpi"]["board"]))
    if gov.female_board_percent is not None:
        gov_kpis.append(("people", f"{gov.female_board_percent:.0f}%", t["kpi"]["women_board"]))
    if gov.independent_board_percent is not None:
        gov_kpis.append(("badge", f"{gov.independent_board_percent:.0f}%", t["kpi"]["independent"]))
    if gov.esg_audit_conducted is not None:
        gov_kpis.append(("badge", t["kpi"]["yes"] if gov.esg_audit_conducted else t["kpi"]["no"], t["kpi"]["audit"]))
    if gov.data_breaches is not None:
        gov_kpis.append(("lock", f"{gov.data_breaches}", t["kpi"]["breaches"]))
    if gov.ethics_violations is not None:
        gov_kpis.append(("scale", f"{gov.ethics_violations}", t["kpi"]["ethics"]))
    if gov.csr_budget_eur is not None:
        gov_kpis.append(("chart", f"{gov.csr_budget_eur:,.0f} €", t["kpi"]["csr_budget"]))
    if gov.sustainability_committee is not None:
        gov_kpis.append(("columns", t["kpi"]["yes"] if gov.sustainability_committee else t["kpi"]["no"], t["kpi"]["committee"]))
    if gov.corruption_cases is not None:
        gov_kpis.append(("scale", f"{gov.corruption_cases}", t["kpi"]["corruption"]))
    pillar_infographic(prs, blank_layout, theme, style, "gov",
                       t["pillar_gov"], _headlines["gov"],
                       scores.governance_score, gov_kpis, t, _insights["gov"])

    # ── SLIDE 5b: Double matérialité (CSRD/ESRS) ────────────────────────
    if "materiality" in chart_images:
        dark = style["dark_slides"]
        slide = content_slide(prs, blank_layout, theme, style,
                              _sh["materiality"], header_color, kicker=t["materiality_title"])
        # Graphique (numéros + légende) à gauche
        add_image_from_bytes(slide, chart_images["materiality"],
                             Inches(0.35), Inches(1.75), height=Inches(4.9))
        # Carte « lecture métier » à droite
        cx, cw = Inches(9.35), Inches(3.65)
        add_shape(slide, ROUNDED_RECT, cx, Inches(2.0), cw, Inches(4.4),
                  fill=theme["card_bg"], line_color=theme["accent"], line_width_pt=1.25)
        add_shape(slide, RECT, cx, Inches(2.0), Inches(0.09), Inches(4.4), fill=theme["accent"])
        add_text(slide, t["key_takeaway"], cx + Inches(0.3), Inches(2.2), cw - Inches(0.5), Inches(0.4),
                 font_size=12, bold=True, color=theme["accent"] if not dark else theme["text_dark"], font=fb)
        add_text(slide, t["materiality_desc"], cx + Inches(0.3), Inches(2.7), cw - Inches(0.55), Inches(3.5),
                 font_size=13.5, color=theme["text_dark"], font=fb)

    # ── SLIDE 5c: Objectifs & trajectoire (hero : chiffre clé + viz) ────
    if "targets" in chart_images or "carbon_trajectory" in chart_images:
        slide = content_slide(prs, blank_layout, theme, style,
                              _sh["objectives"], header_color, kicker=t["targets_title"])
        co2 = request.environmental.co2_emissions_tonnes
        has_carbon = "carbon_trajectory" in chart_images and co2 and co2 > 0
        # Chiffre clé dominant à gauche
        if has_carbon:
            big, cap = "-42%", t["obj_kpi_cap"]
            sub = f"{t['obj_target']} 2030 : {co2 * 0.58:,.0f} {t['obj_unit']}"
        else:
            avg_target = round((min(100, scores.environmental_score + 15) +
                                min(100, scores.social_score + 10) +
                                min(100, scores.governance_score + 5)) / 3)
            big, cap = f"{avg_target}", t["pillars_caption"]
            sub = f"{t['obj_target']} {max(request.company.target_year, request.company.reporting_year + 1)}"
        add_text(slide, big, Inches(0.45), Inches(1.7), Inches(5.4), Inches(1.9),
                 font_size=96, bold=True, color=theme["accent"], font=ft)
        add_text(slide, cap, Inches(0.6), Inches(3.75), Inches(5.3), Inches(0.9),
                 font_size=15, bold=True, color=theme["text_dark"], font=fb)
        add_text(slide, sub, Inches(0.6), Inches(4.55), Inches(5.3), Inches(0.5),
                 font_size=13, color=theme["muted"], font=fb)
        # Carte lecture métier
        add_shape(slide, ROUNDED_RECT, Inches(0.55), Inches(5.15), Inches(5.35), Inches(1.85),
                  fill=theme["card_bg"], line_color=theme["accent"], line_width_pt=1.25)
        add_shape(slide, RECT, Inches(0.55), Inches(5.15), Inches(0.09), Inches(1.85), fill=theme["accent"])
        add_text(slide, t["key_takeaway"], Inches(0.85), Inches(5.32), Inches(4.8), Inches(0.4),
                 font_size=11.5, bold=True, color=theme["accent"], font=fb)
        add_text(slide, t["obj_insight"], Inches(0.85), Inches(5.72), Inches(4.85), Inches(1.2),
                 font_size=13, color=theme["text_dark"], font=fb)
        # Visualisation à droite
        chart_key = "carbon_trajectory" if has_carbon else "targets"
        add_image_from_bytes(slide, chart_images[chart_key],
                             Inches(6.35), Inches(2.3), width=Inches(6.6))

    # ── SLIDE 5d: Taxonomie UE (hero : % dominant + viz) ────────────────
    if "taxonomy" in chart_images:
        slide = content_slide(prs, blank_layout, theme, style,
                              _sh["taxonomy"], header_color, kicker=t["taxonomy_title"])
        tx = request.taxonomy
        vals = [v for v in (tx.turnover_aligned_percent, tx.capex_aligned_percent,
                            tx.opex_aligned_percent) if v is not None]
        top = max(vals) if vals else 0
        add_text(slide, f"{top:.0f}%", Inches(0.45), Inches(1.7), Inches(5.4), Inches(1.9),
                 font_size=96, bold=True, color=theme["env"], font=ft)
        add_text(slide, t["tax_kpi_cap"], Inches(0.6), Inches(3.75), Inches(5.3), Inches(0.9),
                 font_size=15, bold=True, color=theme["text_dark"], font=fb)
        add_shape(slide, ROUNDED_RECT, Inches(0.55), Inches(4.7), Inches(5.35), Inches(2.3),
                  fill=theme["card_bg"], line_color=theme["env"], line_width_pt=1.25)
        add_shape(slide, RECT, Inches(0.55), Inches(4.7), Inches(0.09), Inches(2.3), fill=theme["env"])
        add_text(slide, t["key_takeaway"], Inches(0.85), Inches(4.88), Inches(4.8), Inches(0.4),
                 font_size=11.5, bold=True, color=theme["env"], font=fb)
        add_text(slide, t["tax_insight"], Inches(0.85), Inches(5.28), Inches(4.85), Inches(1.6),
                 font_size=13, color=theme["text_dark"], font=fb)
        add_image_from_bytes(slide, chart_images["taxonomy"],
                             Inches(6.3), Inches(2.5), width=Inches(6.7))

    # ── SLIDE 6: Analyse Stratégique — Solide vs. À progresser ──────────
    slide = content_slide(prs, blank_layout, theme, style, _sh["strategic"], header_color, kicker=t["strategic"])
    dark = style["dark_slides"]
    red = RGBColor(0xE7, 0x4C, 0x3C)
    panel_shape = ROUNDED_RECT if style["card"] != "flat" else RECT
    for (px, pw, ptitle, pcolor, items) in [
        (Inches(0.3), Inches(6.2), t["strengths"], theme["env"], scores.strengths),
        (Inches(6.83), Inches(6.2), t["weaknesses"], red, scores.weaknesses),
    ]:
        # Panneau
        add_shape(slide, panel_shape, px, Inches(1.25), pw, Inches(5.75), fill=theme["card_bg"],
                  line_color=pcolor, line_width_pt=1.0)
        # En-tête chiffré : gros compte + libellé
        add_text(slide, f"{len(items)}", px + Inches(0.3), Inches(1.45), Inches(1.3), Inches(1.0),
                 font_size=46, bold=True, color=pcolor, font=ft)
        add_text(slide, ptitle.upper(), px + Inches(1.45), Inches(1.72), pw - Inches(1.6), Inches(0.7),
                 font_size=15, bold=True, color=theme["text_dark"], font=fb)
        add_bg_rect(slide, px + Inches(0.3), Inches(2.65), pw - Inches(0.6), Inches(0.02),
                    pcolor if not dark else theme["muted"])
        # Items : marqueur coloré + texte
        for i, item in enumerate(items[:5]):
            iy = Inches(2.95) + i * Inches(0.78)
            add_shape(slide, OVAL, px + Inches(0.32), iy + Inches(0.08), Inches(0.16), Inches(0.16), fill=pcolor)
            add_text(slide, item, px + Inches(0.68), iy, pw - Inches(0.95), Inches(0.72),
                     font_size=12.5, color=theme["text_dark"], font=fb)

    # ── SLIDE 7: Recommandations (cartes enrichies pleine largeur) ───────
    if request.include_recommendations:
        from content_generator import enriched_recommendations
        recs = enriched_recommendations(request, scores)[:5]
        slide = content_slide(prs, blank_layout, theme, style, t["recommendations"], header_color)
        add_text(slide, t["rec_intro"], Inches(0.6), Inches(1.2), Inches(12.1), Inches(0.5),
                 font_size=13, color=theme["muted"], font=fb)

        n = len(recs)
        top0 = Inches(1.95)
        row_h = min(Inches(1.02), (Inches(5.15) - (n - 1) * Inches(0.14)) / max(1, n))
        panel_shape = ROUNDED_RECT if style["card"] == "rounded" else RECT
        for i, rec in enumerate(recs):
            pcolor = theme[rec["pillar"]]
            y = top0 + i * (row_h + Inches(0.14))
            # carte
            add_shape(slide, panel_shape, Inches(0.6), y, Inches(12.1), row_h,
                      fill=theme["card_bg"], line_color=pcolor, line_width_pt=1.25)
            # pastille numéro colorée
            add_shape(slide, OVAL, Inches(0.85), y + Inches(0.24), Inches(0.55), Inches(0.55), fill=pcolor)
            add_text(slide, str(i + 1), Inches(0.85), y + Inches(0.30), Inches(0.55), Inches(0.45),
                     font_size=18, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER, font=ft)
            # titre + détail
            add_text(slide, rec["title"], Inches(1.65), y + Inches(0.12), Inches(8.7), Inches(0.45),
                     font_size=15, bold=True, color=theme["text_dark"], font=ft)
            add_text(slide, rec["detail"], Inches(1.65), y + Inches(0.55), Inches(9.2), Inches(0.5),
                     font_size=11, color=theme["muted"], font=fb)
            # chip horizon
            chip_w = Inches(1.7)
            add_shape(slide, ROUNDED_RECT, Inches(10.75), y + Inches(0.28), chip_w, Inches(0.46), fill=pcolor)
            add_text(slide, rec["horizon"], Inches(10.75), y + Inches(0.35), chip_w, Inches(0.35),
                     font_size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER, font=fb)

    # ── SLIDE 8: Alignement ODD ──────────────────────────────────────────
    slide = content_slide(prs, blank_layout, theme, style,
                          _sh["odd"], header_color, kicker=t["odd_title"])

    _odd_colors = [theme["accent"], theme["social"], theme["gov"], theme["env"], theme["env"], theme["gov"]]
    odds = [(num, lab, _odd_colors[i]) for i, (num, lab) in enumerate(t["odds"])]
    odd_shape = ROUNDED_RECT if style["card"] == "rounded" else RECT
    for idx, (odd_num, odd_label, odd_color) in enumerate(odds):
        col = Inches(0.3) + (idx % 3) * Inches(4.3)
        row = Inches(1.3) + (idx // 3) * Inches(1.8)
        if style["card"] == "outline":
            add_shape(slide, RECT, col, row, Inches(4.0), Inches(1.6), fill=theme["card_bg"],
                      line_color=odd_color, line_width_pt=1.5)
            txt_color, sub_color = odd_color, theme["text_dark"]
        else:
            add_shape(slide, odd_shape, col, row, Inches(4.0), Inches(1.6), fill=odd_color)
            txt_color = sub_color = RGBColor(0xFF, 0xFF, 0xFF)
        add_text(slide, odd_num, col + Inches(0.15), row + Inches(0.1),
                 Inches(3.7), Inches(0.6), font_size=18, bold=True, color=txt_color, font=ft)
        add_text(slide, odd_label, col + Inches(0.15), row + Inches(0.75),
                 Inches(3.7), Inches(0.7), font_size=13, color=sub_color, font=fb)

    frameworks = ["GRI Standards", "TCFD", "CSRD / DPEF", "SFDR", "ISO 14001 / 26000"]
    sep_color = theme["accent"] if style["header"] != "minimal" else RGBColor(0xBD, 0xBD, 0xBD)
    add_bg_rect(slide, Inches(0.3), Inches(5.0), Inches(12.73), Inches(0.03), sep_color)
    add_text(slide, t["frameworks_line"] + "  |  ".join(frameworks),
             Inches(0.3), Inches(5.15), Inches(12.7), Inches(0.5),
             font_size=13, bold=True, color=theme["text_dark"], align=PP_ALIGN.CENTER, font=fb)
    # Encart « Impact business »
    _dark = style["dark_slides"]
    add_shape(slide, ROUNDED_RECT, Inches(0.3), Inches(5.85), Inches(12.73), Inches(1.25),
              fill=theme["card_bg"], line_color=theme["accent"], line_width_pt=1.25)
    add_shape(slide, RECT, Inches(0.3), Inches(5.85), Inches(0.09), Inches(1.25), fill=theme["accent"])
    add_text(slide, t["key_takeaway"], Inches(0.6), Inches(6.02), Inches(12), Inches(0.35),
             font_size=12, bold=True, color=theme["accent"] if not _dark else theme["text_dark"], font=fb)
    add_text(slide, t["odd_insight"], Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.65),
             font_size=13.5, color=theme["text_dark"], font=fb)

    # ── SLIDE 9: Conclusion — clôture éditoriale + engagements ──────────
    from content_generator import enriched_recommendations
    _eng = enriched_recommendations(request, scores)[:3]
    slide = prs.slides.add_slide(blank_layout)
    add_bg_rect(slide, 0, 0, SLIDE_W, SLIDE_H, theme["bg_primary"])
    add_bg_rect(slide, 0, 0, Inches(0.18), SLIDE_H, theme["accent"])
    tl_color, sub_color = theme["text_light"], theme["subtitle"]

    add_text(slide, maybe_upper(t["conclusion_title"], style), Inches(0.65), Inches(0.7),
             Inches(11.5), Inches(0.9), font_size=30, bold=True, color=tl_color, font=ft)
    add_bg_rect(slide, Inches(0.68), Inches(1.62), Inches(1.8), Inches(0.06), theme["accent"])
    # Synthèse en grand (l'idée qui reste)
    add_text(slide, _verdict, Inches(0.65), Inches(1.95), Inches(12.0), Inches(1.3),
             font_size=20, bold=True, color=tl_color, font=ft)

    # Engagements prioritaires (les prochaines étapes)
    add_text(slide, t["commitments"], Inches(0.65), Inches(3.55), Inches(8), Inches(0.4),
             font_size=13, bold=True, color=theme["accent"], font=fb)
    ecolors = [theme["env"], theme["social"], theme["gov"]]
    for i, eng in enumerate(_eng):
        y = Inches(4.15) + i * Inches(0.86)
        add_text(slide, f"{i + 1:02d}", Inches(0.65), y, Inches(0.9), Inches(0.7),
                 font_size=30, bold=True, color=theme["accent"], font=ft)
        add_text(slide, eng["title"], Inches(1.7), y + Inches(0.05), Inches(9.0), Inches(0.5),
                 font_size=16, bold=True, color=tl_color, font=fb)
        add_text(slide, eng["horizon"], Inches(10.9), y + Inches(0.08), Inches(1.9), Inches(0.4),
                 font_size=13, bold=True, color=sub_color, align=PP_ALIGN.RIGHT, font=fb)
        if i < len(_eng) - 1:
            add_bg_rect(slide, Inches(1.7), y + Inches(0.72), Inches(11.1), Inches(0.012),
                        RGBColor(0x55, 0x5B, 0x6B) if not style["dark_slides"] else theme["card_bg"])

    add_text(slide, f"© {request.company.reporting_year} {request.company.name} — " + t["confidential"],
             Inches(0.65), Inches(7.0), Inches(11), Inches(0.4),
             font_size=9, italic=True, color=sub_color, font=fb)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
